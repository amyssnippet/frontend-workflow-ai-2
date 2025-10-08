import subprocess
import os
import requests
import json
import re
import mimetypes
import base64
import io
from typing import List, Dict, Tuple
import asyncio
from concurrent.futures import ThreadPoolExecutor, ProcessPoolExecutor
from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Body, Request
from fastapi.responses import JSONResponse, FileResponse, HTMLResponse
import mammoth
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from PIL import Image, ImageDraw, ImageFont

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
import sqlite3
import jwt
from datetime import datetime, timedelta
import hashlib
import uuid
import time
import secrets
from typing import Optional

app = FastAPI(title="Merged Ollama API (llm3)")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://92.242.187.224:3001", "https://flowai.othersys.com"],
    allow_methods=["*"],
    allow_headers=["*"],
)

if not os.path.exists("static"):
    os.makedirs("static")
if not os.path.exists("static/process_images"):
    os.makedirs("static/process_images")
if not os.path.exists("input/images"):
    os.makedirs("input/images")

app.mount("/static", StaticFiles(directory="static"), name="static")

# Ensure data directory and DB
DB_PATH = os.path.join("data", "app.db")
os.makedirs(os.path.dirname(DB_PATH), exist_ok=True)


def get_db_conn():
    conn = sqlite3.connect(DB_PATH, check_same_thread=False)
    conn.row_factory = sqlite3.Row
    return conn


def init_db():
    conn = get_db_conn()
    cur = conn.cursor()
    cur.execute("""
    CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY,
        username TEXT UNIQUE,
        email TEXT UNIQUE,
        password_hash TEXT,
        password_salt TEXT,
        token TEXT,
        token_expiry INTEGER
    )
    """)
    cur.execute("""
    CREATE TABLE IF NOT EXISTS history (
        id INTEGER PRIMARY KEY,
        user_id INTEGER,
        type TEXT,
        title TEXT,
        content TEXT,
        created_at INTEGER,
        updated_at INTEGER,
        FOREIGN KEY(user_id) REFERENCES users(id)
    )
    """)
    cur.execute("""
    CREATE TABLE IF NOT EXISTS sessions (
        id INTEGER PRIMARY KEY,
        user_id INTEGER,
        refresh_token TEXT,
        expires_at INTEGER,
        created_at INTEGER,
        FOREIGN KEY(user_id) REFERENCES users(id)
    )
    """)
    conn.commit()
    conn.close()
    # Attempt to add email verification columns to users table if they don't exist
    try:
        conn = get_db_conn()
        cur = conn.cursor()
        cur.execute("ALTER TABLE users ADD COLUMN email_verified INTEGER DEFAULT 0")
        cur.execute("ALTER TABLE users ADD COLUMN email_verification_token TEXT")
        cur.execute("ALTER TABLE users ADD COLUMN email_verification_expiry INTEGER")
        conn.commit()
        conn.close()
    except Exception:
        # columns already exist or cannot be added; ignore
        try:
            conn.close()
        except:
            pass


# Hardcoded configuration (development). Replace these placeholders with real values if needed.
SMTP_HOST = 'smtp.hostinger.com'
SMTP_PORT = 587
SMTP_USER = 'noreply@othersys.com'  # <-- replace with SMTP username
SMTP_PASS = 'Shidhi@4567'          # <-- replace with SMTP password or app-specific password
EMAIL_FROM = SMTP_USER
FRONTEND_URL = 'https://flowai.othersys.com'

# Development toggle: allow sign-in even if email is not verified
# Hardcoded to True for convenience in local dev. Set to False in production.
ALLOW_UNVERIFIED_LOGIN = True

import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart


def create_email_verification_token(user_id: int, hours_valid: int = 24) -> str:
    token = uuid.uuid4().hex
    expiry = int(time.time()) + hours_valid * 3600
    conn = get_db_conn()
    cur = conn.cursor()
    cur.execute('UPDATE users SET email_verification_token = ?, email_verification_expiry = ?, email_verified = 0 WHERE id = ?', (token, expiry, user_id))
    conn.commit()
    conn.close()
    return token


def send_verification_email(to_email: str, token: str, username: Optional[str] = None):
    if not SMTP_USER or not SMTP_PASS:
        # SMTP not configured — do not send email. Return False so callers can handle.
        return False
    verify_url = f"{FRONTEND_URL}/verify?token={token}"
    subject = "Verify your email"
    text = f"Hello {username or ''},\n\nPlease verify your email by clicking the link below:\n\n{verify_url}\n\nThis link will expire in 24 hours.\n\nThanks!"
    html = f"<p>Hello {username or ''},</p><p>Please verify your email by clicking the link below:</p><p><a href=\"{verify_url}\">Verify Email</a></p><p>This link will expire in 24 hours.</p>"
    msg = MIMEMultipart('alternative')
    msg['Subject'] = subject
    msg['From'] = EMAIL_FROM
    msg['To'] = to_email
    part1 = MIMEText(text, 'plain')
    part2 = MIMEText(html, 'html')
    msg.attach(part1)
    msg.attach(part2)
    try:
        server = smtplib.SMTP(SMTP_HOST, SMTP_PORT)
        server.ehlo()
        if SMTP_PORT == 587:
            server.starttls()
            server.ehlo()
        server.login(SMTP_USER, SMTP_PASS)
        server.sendmail(EMAIL_FROM, [to_email], msg.as_string())
        server.quit()
        return True
    except Exception as e:
        print('Failed to send verification email:', e)
        return False


def hash_password(password: str) -> tuple:
    salt = secrets.token_bytes(16)
    dk = hashlib.pbkdf2_hmac('sha256', password.encode('utf-8'), salt, 100_000)
    return dk.hex(), salt.hex()


def verify_password(password: str, hash_hex: str, salt_hex: str) -> bool:
    salt = bytes.fromhex(salt_hex)
    dk = hashlib.pbkdf2_hmac('sha256', password.encode('utf-8'), salt, 100_000)
    return secrets.compare_digest(dk.hex(), hash_hex)


def create_user(username: Optional[str], email: Optional[str], password: str):
    conn = get_db_conn()
    cur = conn.cursor()
    password_hash, password_salt = hash_password(password)
    try:
        cur.execute("INSERT INTO users (username, email, password_hash, password_salt) VALUES (?,?,?,?)",
                    (username, email, password_hash, password_salt))
        conn.commit()
        user_id = cur.lastrowid
    except sqlite3.IntegrityError as e:
        conn.close()
        raise
    conn.close()
    return user_id


def authenticate_user(identifier: str, password: str):
    conn = get_db_conn()
    cur = conn.cursor()
    cur.execute("SELECT * FROM users WHERE username = ? OR email = ?", (identifier, identifier))
    row = cur.fetchone()
    conn.close()
    if not row:
        return None
    if verify_password(password, row['password_hash'], row['password_salt']):
        return dict(row)
    return None


def generate_token_for_user(user_id: int, days_valid: int = 30) -> str:
    # legacy token storage retained for compatibility; prefer refresh tokens below
    token = uuid.uuid4().hex
    expiry = int(time.time()) + days_valid * 24 * 3600
    conn = get_db_conn()
    cur = conn.cursor()
    cur.execute("UPDATE users SET token = ?, token_expiry = ? WHERE id = ?", (token, expiry, user_id))
    conn.commit()
    conn.close()
    return token


# JWT / Refresh token configuration
# Hardcoded secret key for JWT (development only). Replace for production.
SECRET_KEY = 'dev-secret-key'
ALGORITHM = 'HS256'
ACCESS_TOKEN_EXPIRE_MINUTES = 15
REFRESH_TOKEN_EXPIRE_DAYS = 30


def create_refresh_token_for_user(user_id: int, days: int = REFRESH_TOKEN_EXPIRE_DAYS) -> str:
    refresh = uuid.uuid4().hex
    expires_at = int(time.time()) + days * 24 * 3600
    now = int(time.time())
    conn = get_db_conn()
    cur = conn.cursor()
    cur.execute('INSERT INTO sessions (user_id, refresh_token, expires_at, created_at) VALUES (?,?,?,?)', (user_id, refresh, expires_at, now))
    conn.commit()
    conn.close()
    return refresh


def verify_refresh_token(refresh_token: str):
    conn = get_db_conn()
    cur = conn.cursor()
    now = int(time.time())
    cur.execute('SELECT * FROM sessions WHERE refresh_token = ? AND expires_at > ?', (refresh_token, now))
    row = cur.fetchone()
    conn.close()
    return dict(row) if row else None


def delete_refresh_token(refresh_token: str):
    conn = get_db_conn()
    cur = conn.cursor()
    cur.execute('DELETE FROM sessions WHERE refresh_token = ?', (refresh_token,))
    conn.commit()
    conn.close()


def delete_all_sessions_for_user(user_id: int):
    conn = get_db_conn()
    cur = conn.cursor()
    cur.execute('DELETE FROM sessions WHERE user_id = ?', (user_id,))
    conn.commit()
    conn.close()


def create_access_token(user_id: int, expires_delta: timedelta | None = None):
    to_encode = {"sub": str(user_id)}
    if expires_delta:
        expire = datetime.utcnow() + expires_delta
    else:
        expire = datetime.utcnow() + timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    to_encode.update({"exp": expire})
    encoded_jwt = jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)
    return encoded_jwt


def decode_access_token(token: str):
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        uid = int(payload.get('sub'))
        return uid
    except Exception:
        return None


def get_user_from_request(request: Request):
    # First try Authorization header as Bearer access token
    auth = request.headers.get('authorization') or request.headers.get('Authorization')
    if auth and auth.lower().startswith('bearer '):
        token = auth.split(None, 1)[1]
        # try JWT
        uid = decode_access_token(token)
        if uid:
            conn = get_db_conn()
            cur = conn.cursor()
            cur.execute('SELECT * FROM users WHERE id = ?', (uid,))
            row = cur.fetchone()
            conn.close()
            return dict(row) if row else None
        # fallback to legacy token stored in users table
        user = get_user_by_token(token)
        return user
    return None


@app.middleware('http')
async def attach_user_middleware(request: Request, call_next):
    try:
        user = get_user_from_request(request)
        request.state.user = user
    except Exception:
        request.state.user = None
    response = await call_next(request)
    return response


def get_user_by_token(token: str):
    if not token:
        return None
    conn = get_db_conn()
    cur = conn.cursor()
    now = int(time.time())
    cur.execute("SELECT * FROM users WHERE token = ? AND token_expiry > ?", (token, now))
    row = cur.fetchone()
    conn.close()
    return dict(row) if row else None


def invalidate_token(token: str):
    conn = get_db_conn()
    cur = conn.cursor()
    cur.execute("UPDATE users SET token = NULL, token_expiry = NULL WHERE token = ?", (token,))
    conn.commit()
    conn.close()


def save_history(user_id: int, htype: str, title: str, content: str) -> int:
    now = int(time.time())
    conn = get_db_conn()
    cur = conn.cursor()
    cur.execute("INSERT INTO history (user_id, type, title, content, created_at, updated_at) VALUES (?,?,?,?,?,?)",
                (user_id, htype, title, content, now, now))
    conn.commit()
    hid = cur.lastrowid
    conn.close()
    return hid


def update_history_entry(entry_id: int, user_id: int, title: Optional[str], content: Optional[str]) -> bool:
    parts = []
    params = []
    if title is not None:
        parts.append("title = ?")
        params.append(title)
    if content is not None:
        parts.append("content = ?")
        params.append(content)
    if not parts:
        return False
    params.extend([int(time.time()), entry_id, user_id])
    conn = get_db_conn()
    cur = conn.cursor()
    cur.execute(f"UPDATE history SET {', '.join(parts)}, updated_at = ? WHERE id = ? AND user_id = ?", params)
    conn.commit()
    changed = cur.rowcount
    conn.close()
    return changed > 0


def list_history_for_user(user_id: int):
    conn = get_db_conn()
    cur = conn.cursor()
    cur.execute("SELECT id, type, title, substr(content,1,200) as snippet, created_at, updated_at FROM history WHERE user_id = ? ORDER BY updated_at DESC", (user_id,))
    rows = cur.fetchall()
    conn.close()
    return [dict(r) for r in rows]


def get_history_entry(entry_id: int, user_id: int):
    conn = get_db_conn()
    cur = conn.cursor()
    cur.execute("SELECT * FROM history WHERE id = ? AND user_id = ?", (entry_id, user_id))
    row = cur.fetchone()
    conn.close()
    return dict(row) if row else None


@app.on_event("startup")
async def _init_db_event():
    init_db()

# Shared configuration
OLLAMA_API_URL = "https://climb-buf-pregnant-arrange.trycloudflare.com/api/chat"
OLLAMA_MODEL = "granite3.3:8b"
VISION_MODEL = "qwen2.5vl:3b"
token = "76042dfac098082e9e59a2841542b7588cfb40d5201b16ddd0c32acc85e627d1"

# Concurrency controls
MAX_CONCURRENT_REQUESTS = 10
IO_WORKERS = 20
CPU_WORKERS = 10


@app.on_event("startup")
async def startup_event():
    app.state.io_executor = ThreadPoolExecutor(max_workers=IO_WORKERS)
    app.state.cpu_executor = ProcessPoolExecutor(max_workers=CPU_WORKERS)
    app.state.semaphore = asyncio.Semaphore(MAX_CONCURRENT_REQUESTS)


@app.on_event("shutdown")
async def shutdown_event():
    try:
        app.state.io_executor.shutdown(wait=False)
    except Exception:
        pass
    try:
        app.state.cpu_executor.shutdown(wait=False)
    except Exception:
        pass


def check_mmdc_installed():
    try:
        mmdc_path = r"/root/.volta/bin/mmdc"
        subprocess.run([mmdc_path, "--version"], check=True, capture_output=True)
        return True
    except (FileNotFoundError, subprocess.CalledProcessError):
        print("Error: Mermaid CLI (mmdc) not found.")
        return False


def sanitize_mermaid_code(mermaid_code: str) -> str:
    mermaid_code = re.sub(r"^```mermaid\s*\n", "", mermaid_code, flags=re.MULTILINE)
    mermaid_code = re.sub(r"```$", "", mermaid_code, flags=re.MULTILINE)

    lines = mermaid_code.splitlines()
    cleaned = []
    has_graph = False
    node_counter = 0

    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue

        if stripped.startswith("graph "):
            has_graph = True
            if not stripped.endswith(";"):
                stripped += ";"
            cleaned.append(stripped)
            continue

        if any(edge in stripped for edge in ["-->", "-.->", "---"]):
            safe_line = stripped.replace('"', "'")
            safe_line = safe_line.replace("{", "[").replace("}", "]")
            cleaned.append(safe_line)
            continue

        match = re.match(r"^(\w+)\[(.*)\]$", stripped)
        if match:
            node_id, label = match.groups()
            safe_label = label.replace('"', "'")
            safe_label = re.sub(r"[^\w\s\-\.,()']", "", safe_label)
            safe_label = " ".join(safe_label.split())
            cleaned.append(f"{node_id}[{safe_label}]")
            continue

        node_counter += 1
        safe_label = stripped.replace('"', "'")
        safe_label = re.sub(r"[^\w\s\-\.,()']", "", safe_label)
        safe_label = " ".join(safe_label.split())
        cleaned.append(f"N{node_counter}[{safe_label}]")

    if not has_graph:
        cleaned.insert(0, "graph TD;")

    return "\n".join(cleaned) + "\n"


def call_ollama_granite(user_prompt):
    system_message_content = (
        "You are ONLY to output a valid Mermaid flowchart code block. "
        "The output MUST be ONLY the Mermaid code block, enclosed in triple backticks with 'mermaid'. "
        "Rules to follow strictly:\n"
        "1. The diagram MUST start with 'graph TD;' or 'graph LR;'.\n"
        "2. Node text inside [] or {} MUST NOT contain parentheses (), commas, colons, semicolons, or special symbols. "
        "Use simple words only.\n"
        "3. If multiple options are needed, represent them as separate nodes or as edge labels, not inside a single node.\n"
        "4. Only output nodes and edges — no explanations, no comments.\n"
        "If input cannot be converted, output a minimal valid diagram: "
        "```mermaid\ngraph TD; A[Invalid Input];\n```"
    )

    messages = [
        {"role": "control", "content": "thinking"},
        {"role": "system", "content": system_message_content},
        {"role": "user", "content": user_prompt},
    ]
    payload = {"model": OLLAMA_MODEL, "messages": messages, "stream": False}
    headers = {"Content-Type": "application/json", "Authorization": f"Bearer {token}"}

    try:
        response = requests.post(OLLAMA_API_URL, headers=headers, data=json.dumps(payload), timeout=120)
        response.raise_for_status()
        result = response.json()
        generated_content = result.get("message", {}).get("content", "").strip()

        mermaid_match = re.search(r"```mermaid\s*([\s\S]*?)```", generated_content)
        if mermaid_match:
            return mermaid_match.group(1).strip()
        else:
            return "graph TD;\nA[No valid Mermaid diagram generated];"
    except Exception as e:
        print(f"Error calling Ollama: {e}")
        return "graph TD;\nA[Error generating diagram];"


def repair_mermaid_with_ollama(broken_code: str) -> str:
    prompt = f"The following Mermaid code is invalid. Please fix and return ONLY valid Mermaid code:\n\n```mermaid\n{broken_code}\n```"
    return call_ollama_granite(prompt)


def translate_mermaid_to_image(mermaid_code: str, output_filename: str, output_format="png") -> tuple:
    if not check_mmdc_installed():
        return False, "Mermaid CLI not found"

    temp_mermaid_file = "temp.mmd"
    output_path = f"static/{output_filename}.{output_format}"

    try:
        with open(temp_mermaid_file, "w", encoding="utf-8") as f:
            f.write(mermaid_code)

        mmdc_path = r"/root/.volta/bin/mmdc"

        subprocess.run([mmdc_path, "-i", temp_mermaid_file, "-o", output_path], check=True)
        return True, output_path
    except subprocess.CalledProcessError as e:
        return False, f"mmdc parse error: {e}"
    except Exception as e:
        return False, str(e)
    finally:
        if os.path.exists(temp_mermaid_file):
            os.remove(temp_mermaid_file)


def extract_text_from_file(file_path: str, mime_type: str) -> str:
    # Reuse the advanced extractor from llm.py
    def call_vision_model(image_data: bytes) -> str:
        try:
            if len(image_data) == 0:
                return ""
            encoded_image = base64.b64encode(image_data).decode('utf-8')
            messages = [
                {
                    "role": "user",
                    "content": "Extract all text from this document image. Include headers, body text, lists, tables, and any other visible text. Preserve the structure and formatting. Return only the extracted text without explanations.",
                    "images": [encoded_image]
                }
            ]
            payload = {"model": VISION_MODEL, "messages": messages, "stream": False, "options": {"temperature": 0.1, "num_ctx": 8192}}
            headers = {"Content-Type": "application/json", "Authorization": f"Bearer {token}"}
            response = requests.post(OLLAMA_API_URL, headers=headers, data=json.dumps(payload), timeout=180)
            if response.status_code == 200:
                result = response.json()
                return result.get("message", {}).get("content", "").strip()
            return ""
        except Exception as e:
            print(f"Vision model error: {e}")
            return ""

    def render_docx_to_image(docx_path: str) -> bytes:
        try:
            doc = DocxDocument(docx_path)
            estimated_lines = len(doc.paragraphs) + sum(len(t.rows) for t in doc.tables)
            img_height = max(1200, estimated_lines * 40 + 100)
            img = Image.new('RGB', (1200, img_height), 'white')
            draw = ImageDraw.Draw(img)
            try:
                font_normal = ImageFont.truetype("arial.ttf", 24)
            except:
                font_normal = ImageFont.load_default()
            y_position = 30
            x_margin = 30
            for para in doc.paragraphs:
                if para.text.strip():
                    text = para.text[:150]
                    draw.text((x_margin, y_position), text, fill='black', font=font_normal)
                    y_position += 35
            for table in doc.tables:
                y_position += 10
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip()[:20] for cell in row.cells if cell.text.strip()])
                    if row_text:
                        draw.text((x_margin, y_position), row_text, fill='black', font=font_normal)
                        y_position += 30
                y_position += 10
            if y_position < img_height:
                img = img.crop((0, 0, 1200, min(y_position + 50, img_height)))
            img_byte_arr = io.BytesIO()
            img.save(img_byte_arr, format='JPEG', quality=70, optimize=True)
            return img_byte_arr.getvalue()
        except Exception as e:
            print(f"Error rendering DOCX: {e}")
            return b""

    text = ""
    try:
        if mime_type == "application/pdf":
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text += f"{page_text}\n"
                    else:
                        page_image = page.to_image(resolution=120)
                        img_byte_arr = io.BytesIO()
                        page_image.original.save(img_byte_arr, format='PNG')
                        vision_text = call_vision_model(img_byte_arr.getvalue())
                        if vision_text:
                            text += f"{vision_text}\n"
        elif mime_type in [
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        ]:
            doc = DocxDocument(file_path)
            paragraphs_text = [p.text for p in doc.paragraphs if p.text.strip()]
            tables_text = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                    if row_text:
                        tables_text.append(row_text)
            all_text = paragraphs_text + tables_text
            text = "\n".join(all_text)
            if not text.strip():
                rendered_image = render_docx_to_image(file_path)
                if rendered_image:
                    text = call_vision_model(rendered_image)
        elif mime_type == "text/plain":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        elif mime_type in ["image/png", "image/jpeg", "image/jpg", "image/bmp", "image/gif"]:
            with open(file_path, "rb") as f:
                text = call_vision_model(f.read())
        else:
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
            except:
                text = ""
    except Exception as e:
        print(f"Extraction failed: {e}")
        return "Error: Unable to process this document format."

    result = text.strip()
    if not result:
        return "This document appears to be empty or contains only images/formatting."
    return result


def call_ollama_vision(image_data: bytes, prompt: str) -> str:
    try:
        encoded_image = base64.b64encode(image_data).decode('utf-8')
        messages = [{"role": "user", "content": prompt, "images": [encoded_image]}]
        payload = {"model": VISION_MODEL, "messages": messages, "stream": False, "options": {"temperature": 0.1, "num_ctx": 4096}}
        headers = {"Content-Type": "application/json", "Authorization": f"Bearer {token}"}
        response = requests.post(OLLAMA_API_URL, headers=headers, data=json.dumps(payload), timeout=120)
        if response.status_code == 200:
            result = response.json()
            return result.get("message", {}).get("content", "").strip()
        else:
            print(f"[VISION] Error response: {response.text[:500]}")
            return ""
    except Exception as e:
        print(f"[VISION] Exception: {e}")
        return ""


def create_process_image(step_number: int, title: str, description: str, status: str, data: Dict = None) -> str:
    width = 800
    height = 400
    img = Image.new('RGB', (width, height), 'white')
    draw = ImageDraw.Draw(img)
    try:
        font_title = ImageFont.truetype("arial.ttf", 32)
        font_desc = ImageFont.truetype("arial.ttf", 20)
        font_number = ImageFont.truetype("arialbd.ttf", 80)
        font_status = ImageFont.truetype("arial.ttf", 18)
    except:
        font_title = ImageFont.load_default()
        font_desc = ImageFont.load_default()
        font_number = ImageFont.load_default()
        font_status = ImageFont.load_default()
    if status == "success":
        bg_color = (34, 139, 34)
        status_text = "✓ SUCCESS"
    elif status == "processing":
        bg_color = (30, 144, 255)
        status_text = "⟳ PROCESSING"
    elif status == "error":
        bg_color = (220, 20, 60)
        status_text = "✗ ERROR"
    else:
        bg_color = (128, 128, 128)
        status_text = "○ PENDING"
    draw.rectangle([(0, 0), (width, 100)], fill=bg_color)
    circle_center = (80, 50)
    circle_radius = 40
    draw.ellipse([(circle_center[0] - circle_radius, circle_center[1] - circle_radius), (circle_center[0] + circle_radius, circle_center[1] + circle_radius)], fill='white')
    number_text = str(step_number)
    try:
        number_bbox = draw.textbbox((0, 0), number_text, font=font_number)
        number_width = number_bbox[2] - number_bbox[0]
        number_height = number_bbox[3] - number_bbox[1]
    except Exception:
        number_width = 40
        number_height = 40
    draw.text((circle_center[0] - number_width // 2, circle_center[1] - number_height // 2 - 5), number_text, fill=bg_color, font=font_number)
    draw.text((140, 30), title, fill='white', font=font_title)
    try:
        status_bbox = draw.textbbox((0, 0), status_text, font=font_status)
        status_width = status_bbox[2] - status_bbox[0]
    except Exception:
        status_width = 100
    draw.text((width - status_width - 20, 35), status_text, fill='white', font=font_status)
    draw.line([(20, 110), (width - 20, 110)], fill=bg_color, width=3)
    y_offset = 140
    words = description.split()
    lines = []
    current_line = []
    for word in words:
        current_line.append(word)
        test_line = ' '.join(current_line)
        try:
            bbox = draw.textbbox((0, 0), test_line, font=font_desc)
            if bbox[2] - bbox[0] > width - 80:
                if len(current_line) > 1:
                    current_line.pop()
                    lines.append(' '.join(current_line))
                    current_line = [word]
                else:
                    lines.append(test_line)
                    current_line = []
        except Exception:
            pass
    if current_line:
        lines.append(' '.join(current_line))
    for line in lines[:5]:
        draw.text((40, y_offset), line, fill='black', font=font_desc)
        y_offset += 30
    if data:
        y_offset += 20
        draw.line([(40, y_offset), (width - 40, y_offset)], fill='lightgray', width=1)
        y_offset += 15
        for key, value in list(data.items())[:3]:
            text = f"{key}: {value}"
            if len(text) > 80:
                text = text[:77] + "..."
            draw.text((40, y_offset), text, fill='gray', font=font_status)
            y_offset += 25
    draw.rectangle([(0, 0), (width - 1, height - 1)], outline=bg_color, width=3)
    filename = f"process_step_{step_number:02d}.png"
    filepath = f"static/process_images/{filename}"
    img.save(filepath, 'PNG')
    return filename


def extract_steps_from_analysis(image_analyses: List[Dict]) -> List[Dict]:
    steps = []
    numbered_re = re.compile(r'^\s*(\d+)\.\s*(.+)$', re.MULTILINE)
    for idx, analysis_data in enumerate(image_analyses, 1):
        analysis = analysis_data.get('analysis', '') or ''
        objects = []
        # try to extract numbered list lines
        lines = [line.strip() for line in analysis.splitlines() if line.strip()]
        for line in lines:
            m = re.match(r'^\s*\d+\.\s*(.+)', line)
            if m:
                # attempt to split 'ObjectName - Description' by first ' - '
                content = m.group(1).strip()
                if ' - ' in content:
                    name, instr = content.split(' - ', 1)
                elif ':' in content:
                    name, instr = content.split(':', 1)
                else:
                    # fallback: split first dash or take whole line as name
                    parts = re.split(r'\s-\s|\s—\s|\s:\s', content, 1)
                    if len(parts) == 2:
                        name, instr = parts[0].strip(), parts[1].strip()
                    else:
                        name, instr = content, ''
                objects.append({'name': name.strip(), 'instruction': instr.strip()})

        # fallback to summarized action when no numbered objects found
        if not objects:
            summary = summarize_process_details(analysis)
            step_description = summary.get('action') or f"Process from {analysis_data.get('filename', 'image')}"
            # add a single object representing the step description
            objects = [{'name': step_description, 'instruction': summary.get('key_details', [])[:3]}]
        steps.append({
            "number": idx,
            "title": f"Step {idx}",
            "description": objects[0]['name'] if objects else f"Step {idx}",
            "objects": objects,
            "full_analysis": analysis,
            "image_path": analysis_data.get('image_path', ''),
            "filename": analysis_data.get('filename', ''),
            "next_step": idx + 1 if idx < len(image_analyses) else None
        })
    return steps


def summarize_process_details(analysis_text: str) -> Dict[str, str]:
    lines = [line.strip() for line in analysis_text.split('\n') if line.strip()]
    summary = {'action': '', 'input': '', 'output': '', 'key_details': []}
    for line in lines:
        clean_line = re.sub(r'^\d+\.\s*', '', line)
        lower_line = clean_line.lower()
        if any(k in lower_line for k in ['action', 'process', 'shown', 'displays', 'interface']):
            if not summary['action']:
                summary['action'] = clean_line[:150]
        elif any(k in lower_line for k in ['input', 'prerequisite', 'required', 'need']):
            if not summary['input']:
                summary['input'] = clean_line[:150]
        elif any(k in lower_line for k in ['output', 'result', 'generate', 'create']):
            if not summary['output']:
                summary['output'] = clean_line[:150]
        elif any(k in lower_line for k in ['setting', 'configuration', 'option', 'field', 'button', 'value']):
            if len(summary['key_details']) < 3:
                summary['key_details'].append(clean_line[:100])
    if not summary['action'] and lines:
        summary['action'] = lines[0][:150]
    if not summary['key_details']:
        summary['key_details'] = [re.sub(r'^\d+\.\s*', '', l)[:100] for l in lines[1:4]]
    return summary


def create_sop_docx_with_images(workflow_steps: List[Dict], flowchart_path: str = None) -> str:
    doc = DocxDocument()
    # Title block
    title = doc.add_heading('Standard Operating Procedure', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle = doc.add_paragraph('Step-by-Step Workflow Documentation')
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph()

    # Use compact spacing; do not insert page breaks between steps so document is continuous.
    for step in workflow_steps:
        # Insert the step header
        h = doc.add_paragraph()
        h_run = h.add_run(f"Step {step['number']}: {step.get('title', '')}")
        h_run.bold = True
        h_run.font.size = Pt(12)

        # Insert image (if exists) scaled to a reasonable width to avoid pushing to new pages
        if step.get('image_path') and os.path.exists(step['image_path']):
            try:
                img_para = doc.add_paragraph()
                img_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = img_para.add_run()
                # Constrain image to 5.5 inches wide to prevent large page jumps
                run.add_picture(step['image_path'], width=Inches(5.5))
            except Exception as e:
                doc.add_paragraph(f"⚠ Unable to add image: {e}")

        # Add a short bold caption for the image/step
        caption = doc.add_paragraph()
        caption_run = caption.add_run(step.get('description', ''))
        caption_run.bold = True
        caption_run.font.size = Pt(10)

        # Add the objects parsed from analysis as a concise bullet list
        objects = step.get('objects', [])
        if objects:
            obj_heading = doc.add_paragraph()
            obj_heading.add_run('Key objects & actions:').bold = True
            for obj in objects:
                # obj['instruction'] may be a list (fallback) or string
                instr = obj.get('instruction', '')
                if isinstance(instr, list):
                    instr_text = '; '.join([str(x) for x in instr])
                else:
                    instr_text = str(instr)

                # Use a concise bullet: ObjectName — short instruction
                b = doc.add_paragraph(style='List Bullet')
                b_run = b.add_run(f"{obj.get('name', '')} — {instr_text}")
                b_run.font.size = Pt(10)

        # Small spacer between steps; avoid page breaks for continuous flow
        doc.add_paragraph()

    docx_filename = "workflow_sop.docx"
    docx_path = f"static/{docx_filename}"
    doc.save(docx_path)
    return docx_path


@app.post("/generate-process-images/")
async def generate_diagram_with_process_images(request: Request, description: str = Form(None), files: List[UploadFile] = File(None), output_format: str = Form("png")):
    async with app.state.semaphore:
        print("Starting request (generate-process-images)")
        process_dir = "static/process_images"
        for f in os.listdir(process_dir):
            if f.startswith("process_step_"):
                os.remove(os.path.join(process_dir, f))
        process_steps = []
        image_analyses = []
        loop = asyncio.get_event_loop()
        if files and len(files) > 0:
            for idx, file in enumerate(files, 1):
                image_path = f"input/images/{file.filename}"
                with open(image_path, "wb") as fh:
                    content = await file.read()
                    fh.write(content)
                with open(image_path, "rb") as fh:
                    image_data = fh.read()
                analysis_prompt = f"This is Step {idx} of a multi-step workflow process. Analyze the screenshot and extract all visible steps."
                analysis = await loop.run_in_executor(app.state.io_executor, call_ollama_vision, image_data, analysis_prompt)
                image_analyses.append({"step_number": idx, "filename": file.filename, "image_path": image_path, "analysis": analysis or f"Unable to analyze step {idx}", "status": 'success' if analysis else 'error'})
                img_filename = await loop.run_in_executor(app.state.cpu_executor, create_process_image, idx, f"Step {idx}: {file.filename}", f"Analyzed workflow step from {file.filename}", "success" if analysis else "error", {"Image": file.filename, "Analysis Length": f"{len(analysis) if analysis else 0} chars"})
                process_steps.append({"step": idx, "title": f"Step {idx}: {file.filename}", "status": 'success' if analysis else 'error', "image": f"/static/process_images/{img_filename}"})
        else:
            raise HTTPException(status_code=400, detail="No images provided; use files parameter")
        workflow_steps = extract_steps_from_analysis(image_analyses)
        sop_docx_path = await loop.run_in_executor(app.state.cpu_executor, create_sop_docx_with_images, workflow_steps)
        result = {"sop_docx_url": f"/static/{os.path.basename(sop_docx_path)}", "process_steps": process_steps, "workflow_steps": workflow_steps, "total_steps": len(workflow_steps)}
        try:
            auth = request.headers.get('authorization') or request.headers.get('Authorization')
            if auth and auth.lower().startswith('bearer '):
                token = auth.split(None, 1)[1]
                user = get_user_by_token(token)
                if user:
                    # Save docx reference
                    save_history(user['id'], 'docx', f'Processed SOP {int(time.time())}', f"/static/{os.path.basename(sop_docx_path)}")
                    # Save first mermaid if available
                    if workflow_steps and workflow_steps[0].get('full_analysis'):
                        save_history(user['id'], 'mermaid', f'Flow from process {int(time.time())}', workflow_steps[0]['full_analysis'])
        except Exception:
            pass
        return result


@app.post("/generate-docx/")
async def generate_diagram_with_docx(request: Request, description: str = Form(None), file: UploadFile = File(None), output_format: str = Form("png")):
    # Keep implementation similar to llm.py but run heavy parts in executors
    prompt_text = description.strip() if description else None
    steps_log = []
    if file is not None:
        steps_log.append("Step 1: File uploaded - extracting text from document")
        input_dir = "input"
        os.makedirs(input_dir, exist_ok=True)
        file_path = os.path.join(input_dir, file.filename)
        with open(file_path, "wb") as fh:
            fh.write(await file.read())
        mime_type, _ = mimetypes.guess_type(file_path)
        if not mime_type:
            raise HTTPException(status_code=400, detail="Unsupported file type")
        extracted_text = extract_text_from_file(file_path, mime_type)
        if not extracted_text:
            raise HTTPException(status_code=422, detail="File is empty or unsupported")
        prompt_text = extracted_text
        steps_log.append(f"Step 1 Complete: Extracted {len(extracted_text)} characters from {file.filename}")
    else:
        steps_log.append("Step 1: Using provided text description")
    if not prompt_text:
        raise HTTPException(status_code=400, detail="No input provided")
    steps_log.append(f"Step 2: Generating Mermaid flowchart code using {OLLAMA_MODEL}")
    mermaid_code = call_ollama_granite(prompt_text)
    mermaid_code_original = mermaid_code
    mermaid_code = sanitize_mermaid_code(mermaid_code)
    steps_log.append("Step 2 Complete: Mermaid code generated and sanitized")
    steps_log.append("Step 3: Rendering Mermaid code to image")
    success, result = translate_mermaid_to_image(mermaid_code, "generated_flowchart", output_format)
    if not success:
        steps_log.append("Step 3 Failed: Initial render failed, attempting repair")
        repaired_code = repair_mermaid_with_ollama(mermaid_code)
        repaired_code = sanitize_mermaid_code(repaired_code)
        success, result = translate_mermaid_to_image(repaired_code, "generated_flowchart", output_format)
        if success:
            mermaid_code = repaired_code
            steps_log.append("Step 4 Complete: Repair successful, diagram generated")
        else:
            mermaid_code = "graph TD;\nA[Diagram generation failed];"
            result = f"static/generated_flowchart.{output_format}"
            with open(result, "wb") as fh:
                fh.write(b"")
            steps_log.append("Step 4 Failed: Unable to generate valid diagram")
    else:
        steps_log.append("Step 3 Complete: Diagram successfully rendered")
    steps_log.append("Step 5: Creating comprehensive DOCX document")
    # Create DOCX synchronously (small) and offload heavier images if necessary
    doc = DocxDocument()
    title = doc.add_heading('Flowchart Generation Report', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_heading('Processing Steps', level=1)
    for step in steps_log:
        p = doc.add_paragraph(step, style='List Bullet')
        p.paragraph_format.space_after = Pt(6)
    doc.add_paragraph()
    doc.add_heading('Input Content', level=1)
    input_para = doc.add_paragraph()
    input_para.add_run('Source: ').bold = True
    input_para.add_run(f"{file.filename if file else 'Text description'}\n")
    input_para.add_run('Length: ').bold = True
    input_para.add_run(f"{len(prompt_text)} characters\n\n")
    display_text = prompt_text if len(prompt_text) <= 2000 else prompt_text[:2000] + "\n\n[... text truncated for brevity ...]"
    doc.add_paragraph(display_text)
    doc.add_heading('Generated Flowchart Diagram', level=1)
    if os.path.exists(result) and os.path.getsize(result) > 0:
        try:
            # constrain flowchart image width to avoid pushing subsequent content to a new page
            doc.add_picture(result, width=Inches(5.5))
            steps_log.append("Step 5: Flowchart image added to document")
        except Exception as e:
            doc.add_paragraph(f"Error adding image: {str(e)}")
    else:
        doc.add_paragraph("⚠ Flowchart image could not be generated or is empty")
    doc.add_paragraph()
    doc.add_heading('Mermaid Source Code', level=1)
    code_para = doc.add_paragraph()
    code_run = code_para.add_run(mermaid_code)
    code_run.font.name = 'Courier New'
    code_run.font.size = Pt(10)
    code_run.font.color.rgb = RGBColor(0, 0, 128)
    doc.add_heading('Generation Metadata', level=1)
    metadata = [f"Model Used: {OLLAMA_MODEL}", f"Vision Model: {VISION_MODEL}", f"Output Format: {output_format.upper()}", f"Mermaid Code Length: {len(mermaid_code)} characters", f"Repair Required: {'Yes' if not success else 'No'}", f"Total Processing Steps: {len(steps_log)}"]
    for item in metadata:
        doc.add_paragraph(item, style='List Bullet')
    docx_filename = "flowchart_generation_report.docx"
    docx_path = f"static/{docx_filename}"
    doc.save(docx_path)
    # If request includes a valid Authorization token, save history entries for generated content
    try:
        auth = request.headers.get('authorization') or request.headers.get('Authorization')
        if auth and auth.lower().startswith('bearer '):
            token = auth.split(None, 1)[1]
            user = get_user_by_token(token)
            if user:
                # Save mermaid
                save_history(user['id'], 'mermaid', f'Generated flowchart {int(time.time())}', mermaid_code)
                # Save docx reference
                save_history(user['id'], 'docx', f'Generated report {int(time.time())}', f"/static/{docx_filename}")
    except Exception:
        pass
    steps_log.append("Step 5 Complete: DOCX document created successfully")
    return {"mermaid": mermaid_code, "image_url": f"/static/{os.path.basename(result)}", "docx_url": f"/static/{docx_filename}", "message": "Flowchart and detailed report generated successfully", "steps": steps_log, "image_exists": os.path.exists(result), "image_size": os.path.getsize(result) if os.path.exists(result) else 0}


@app.post("/regenerate-docx/")
async def regenerate_docx(workflow_steps: List[Dict] = Body(...)):
    loop = asyncio.get_event_loop()
    formatted_steps = []
    for step in workflow_steps:
        formatted_steps.append({
            "number": step.get("number"),
            "title": step.get("title"),
            "description": step.get("description"),
            "summary": step.get("summary", {}),
            "full_analysis": step.get("full_analysis", ""),
            "image_path": step.get("image_path", ""),
            "filename": step.get("filename", ""),
            "next_step": step.get("next_step")
        })
    docx_path = await loop.run_in_executor(app.state.cpu_executor, create_sop_docx_with_images, formatted_steps)
    return {"sop_docx_url": f"/static/{os.path.basename(docx_path)}", "message": "DOCX regenerated successfully", "total_steps": len(formatted_steps)}


@app.get("/concurrency-status/")
async def concurrency_status():
    sem = app.state.semaphore
    used = MAX_CONCURRENT_REQUESTS - sem._value
    return {"max_concurrent_requests": MAX_CONCURRENT_REQUESTS, "used_slots": used, "available_slots": sem._value}


@app.post("/simulate-parallel/")
async def simulate_parallel(count: int = Form(10), dummy_seconds: int = Form(3)):
    if count < 1:
        raise HTTPException(status_code=400, detail="count must be >= 1")
    async def _run_dummy(i: int):
        loop = asyncio.get_event_loop()
        def cpu_work(n):
            import time
            end = time.time() + n
            s = 0
            while time.time() < end:
                s += sum(i * i for i in range(1000))
            return s
        return await loop.run_in_executor(app.state.cpu_executor, cpu_work, dummy_seconds)
    tasks = [asyncio.create_task(_run_dummy(i)) for i in range(count)]
    results = await asyncio.gather(*tasks)
    return {"ran": count, "results_sample": results[:3]}


@app.post("/generate/")
async def generate_diagram(description: str = Form(None), file: UploadFile = File(None), output_format: str = Form("png")):
    raise HTTPException(status_code=400, detail="Please use /generate-process-images/ or /generate-docx/ for workflow generation")


@app.get("/static/{image_name}")
async def get_diagram_image(image_name: str):
    path = f"static/{image_name}"
    if not os.path.exists(path):
        raise HTTPException(404, "Image not found")
    media_type = "image/svg+xml" if image_name.endswith(".svg") else "image/png"
    return FileResponse(path, media_type=media_type)


@app.get("/preview-docx/{filename}")
def preview_docx(filename: str):
    # Return an HTML preview of a DOCX generated by the server using mammoth
    safe_name = os.path.basename(filename)
    path = os.path.join("static", safe_name)
    if not os.path.exists(path):
        raise HTTPException(404, "DOCX not found")
    try:
        with open(path, "rb") as f:
            result = mammoth.convert_to_html(f)
            html = result.value
        wrapper = (
            "<html><head><meta charset='utf-8'><meta name='viewport' content='width=device-width,initial-scale=1'>"
            "<style>body{font-family:system-ui,-apple-system,Segoe UI,Roboto,'Helvetica Neue',Arial;padding:20px;}</style></head>"
            f"<body>{html}</body></html>"
        )
        return HTMLResponse(content=wrapper, status_code=200)
    except Exception as e:
        print(f"Preview conversion error: {e}")
        raise HTTPException(500, "Failed to convert DOCX to HTML")


@app.get("/")
async def root():
    return {"message": "Merged llm3 API running - endpoints: /generate-process-images/, /generate-docx/, /regenerate-docx/"}


@app.post('/signup')
async def signup(payload: Dict = Body(...)):
    # Accept fullName for frontend compatibility
    username = payload.get('username') or payload.get('fullName')
    email = payload.get('email')
    password = payload.get('password')
    if not password or (not username and not email):
        raise HTTPException(status_code=400, detail='username/email and password required')
    try:
        user_id = create_user(username, email, password)
        # create verification token and send email
        verification_sent = False
        if email:
            token = create_email_verification_token(user_id)
            verification_sent = send_verification_email(email, token, username)
        return {'id': user_id, 'message': 'User created', 'verification_sent': bool(verification_sent)}
    except sqlite3.IntegrityError as e:
        raise HTTPException(status_code=400, detail='User with that username or email already exists')


@app.post('/signin')
async def signin(payload: Dict = Body(...)):
    identifier = payload.get('identifier') or payload.get('username') or payload.get('email')
    password = payload.get('password')
    if not identifier or not password:
        raise HTTPException(status_code=400, detail='identifier and password required')
    user = authenticate_user(identifier, password)
    if not user:
        raise HTTPException(status_code=401, detail='Invalid credentials')
    # block if email not verified unless dev toggle enabled
    if user.get('email') and not user.get('email_verified'):
        if not ALLOW_UNVERIFIED_LOGIN:
            raise HTTPException(status_code=403, detail='Email not verified. Please check your inbox or request a verification email.')
        else:
            # allow login but inform client
            access = create_access_token(user['id'])
            refresh = create_refresh_token_for_user(user['id'])
            generate_token_for_user(user['id'])
            return {'access_token': access, 'refresh_token': refresh, 'user': {'id': user['id'], 'username': user.get('username'), 'email': user.get('email')}, 'warning': 'Email not verified, logged in due to ALLOW_UNVERIFIED_LOGIN'}
    # Create refresh session and issue access token
    refresh = create_refresh_token_for_user(user['id'])
    access = create_access_token(user['id'])
    # keep legacy token as well
    generate_token_for_user(user['id'])
    return {'access_token': access, 'refresh_token': refresh, 'user': {'id': user['id'], 'username': user.get('username'), 'email': user.get('email')}}


@app.post('/signout')
async def signout(request: Request, payload: Dict = Body(None)):
    # Sign out a refresh token (client should send refresh_token) or invalidate legacy token in header
    if payload and payload.get('refresh_token'):
        delete_refresh_token(payload['refresh_token'])
        return {'message': 'Refresh token deleted'}
    auth = request.headers.get('authorization') or request.headers.get('Authorization')
    if not auth or not auth.lower().startswith('bearer '):
        raise HTTPException(status_code=401, detail='Authorization required')
    token = auth.split(None, 1)[1]
    # attempt to delete as refresh token
    try:
        delete_refresh_token(token)
    except Exception:
        pass
    # invalidate legacy token if present
    invalidate_token(token)
    return {'message': 'Signed out'}


@app.post('/refresh')
async def refresh_token(payload: Dict = Body(...)):
    refresh = payload.get('refresh_token')
    if not refresh:
        raise HTTPException(status_code=400, detail='refresh_token required')
    session = verify_refresh_token(refresh)
    if not session:
        raise HTTPException(status_code=401, detail='Invalid or expired refresh token')
    user_id = session['user_id']
    access = create_access_token(user_id)
    return {'access_token': access}


@app.post('/signout-all')
async def signout_all(request: Request):
    user = getattr(request.state, 'user', None)
    if not user:
        raise HTTPException(status_code=401, detail='Authorization required')
    delete_all_sessions_for_user(user['id'])
    return {'message': 'All sessions deleted'}


@app.get('/verify-email')
async def verify_email(token: str = None):
    if not token:
        raise HTTPException(status_code=400, detail='token required')
    conn = get_db_conn()
    cur = conn.cursor()
    now = int(time.time())
    cur.execute('SELECT id, email_verification_expiry FROM users WHERE email_verification_token = ?', (token,))
    row = cur.fetchone()
    if not row:
        conn.close()
        raise HTTPException(status_code=404, detail='Invalid token')
    if row['email_verification_expiry'] and row['email_verification_expiry'] < now:
        conn.close()
        raise HTTPException(status_code=400, detail='Token expired')
    cur.execute('UPDATE users SET email_verified = 1, email_verification_token = NULL, email_verification_expiry = NULL WHERE id = ?', (row['id'],))
    conn.commit()
    conn.close()
    # Redirect to frontend success page
    redirect_url = f"{FRONTEND_URL}/verify-success"
    return JSONResponse({'message': 'Email verified', 'redirect': redirect_url})


@app.post('/resend-verification')
async def resend_verification(payload: Dict = Body(...)):
    identifier = payload.get('identifier') or payload.get('email') or payload.get('username')
    if not identifier:
        raise HTTPException(status_code=400, detail='identifier required')
    conn = get_db_conn()
    cur = conn.cursor()
    cur.execute('SELECT id, email FROM users WHERE username = ? OR email = ?', (identifier, identifier))
    row = cur.fetchone()
    conn.close()
    if not row:
        raise HTTPException(status_code=404, detail='User not found')
    if not row['email']:
        raise HTTPException(status_code=400, detail='User has no email')
    token = create_email_verification_token(row['id'])
    sent = send_verification_email(row['email'], token)
    return {'message': 'Verification email queued' if sent else 'Verification email not sent', 'verification_sent': bool(sent)}


# Compatibility auth-scoped routes (some clients expect /auth/* paths)
@app.post('/auth/signup')
async def auth_signup(payload: Dict = Body(...)):
    return await signup(payload)


@app.post('/auth/signin')
async def auth_signin(payload: Dict = Body(...)):
    return await signin(payload)


@app.post('/auth/resend-verification')
async def auth_resend_verification(payload: Dict = Body(...)):
    return await resend_verification(payload)


@app.get('/auth/verify')
async def auth_verify(token: str = None):
    return await verify_email(token)


@app.get('/auth/me')
async def auth_me(request: Request):
    return await me(request)


@app.post('/auth/signout')
async def auth_signout(request: Request, payload: Dict = Body(None)):
    return await signout(request, payload)


@app.get('/me')
async def me(request: Request):
    # Prefer the user attached by middleware (supports JWT access tokens)
    user = getattr(request.state, 'user', None)
    if user:
        return {'id': user['id'], 'username': user.get('username'), 'email': user.get('email')}

    # Fallback: try to decode Authorization header manually (support legacy token lookup as well)
    auth = request.headers.get('authorization') or request.headers.get('Authorization')
    if auth and auth.lower().startswith('bearer '):
        token = auth.split(None, 1)[1]
        # Try JWT access token first
        uid = decode_access_token(token)
        if uid:
            conn = get_db_conn()
            cur = conn.cursor()
            cur.execute('SELECT id, username, email FROM users WHERE id = ?', (uid,))
            row = cur.fetchone()
            conn.close()
            if row:
                return {'id': row['id'], 'username': row['username'], 'email': row['email']}

        # Fall back to legacy token stored on users table
        user = get_user_by_token(token)
        if user:
            return {'id': user['id'], 'username': user.get('username'), 'email': user.get('email')}

    raise HTTPException(status_code=401, detail='Authorization required or invalid token')


@app.post('/history/save')
async def history_save(request: Request, payload: Dict = Body(...)):
    user = getattr(request.state, 'user', None)
    if not user:
        raise HTTPException(status_code=401, detail='Authorization required')
    htype = payload.get('type')
    title = payload.get('title') or f'{htype} saved {int(time.time())}'
    content = payload.get('content') or ''
    entry_id = save_history(user['id'], htype, title, content)
    return {'id': entry_id}


@app.post('/history/update/{entry_id}')
async def history_update(entry_id: int, request: Request, payload: Dict = Body(...)):
    user = getattr(request.state, 'user', None)
    if not user:
        raise HTTPException(status_code=401, detail='Authorization required')
    title = payload.get('title')
    content = payload.get('content')
    ok = update_history_entry(entry_id, user['id'], title, content)
    if not ok:
        raise HTTPException(status_code=404, detail='Not found or no changes')
    return {'id': entry_id}


@app.get('/history/list')
async def history_list(request: Request):
    user = getattr(request.state, 'user', None)
    if not user:
        raise HTTPException(status_code=401, detail='Authorization required')
    items = list_history_for_user(user['id'])
    return {'items': items}


@app.get('/history/{entry_id}')
async def history_get(entry_id: int, request: Request):
    user = getattr(request.state, 'user', None)
    if not user:
        raise HTTPException(status_code=401, detail='Authorization required')
    entry = get_history_entry(entry_id, user['id'])
    if not entry:
        raise HTTPException(status_code=404, detail='Not found')
    return {'item': entry}


@app.delete('/history/{entry_id}')
async def history_delete(entry_id: int, request: Request):
    user = getattr(request.state, 'user', None)
    if not user:
        raise HTTPException(status_code=401, detail='Authorization required')
    conn = get_db_conn()
    cur = conn.cursor()
    cur.execute('DELETE FROM history WHERE id = ? AND user_id = ?', (entry_id, user['id']))
    conn.commit()
    deleted = cur.rowcount
    conn.close()
    if not deleted:
        raise HTTPException(status_code=404, detail='Not found')
    return {'deleted': True}


@app.post('/autosave')
async def autosave(request: Request, payload: Dict = Body(...)):
    """Autosave an editing buffer. If 'id' provided, update that history entry; otherwise create a new one."""
    user = getattr(request.state, 'user', None)
    if not user:
        raise HTTPException(status_code=401, detail='Authorization required')
    entry_id = payload.get('id')
    htype = payload.get('type', 'mermaid')
    title = payload.get('title')
    content = payload.get('content', '')
    if entry_id:
        ok = update_history_entry(int(entry_id), user['id'], title, content)
        if not ok:
            raise HTTPException(status_code=404, detail='Entry not found')
        return {'id': entry_id, 'updated': True}
    else:
        hid = save_history(user['id'], htype, title or f'Autosave {int(time.time())}', content)
        return {'id': hid, 'created': True}
